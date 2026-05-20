from __future__ import annotations

import re
import shutil
from pathlib import Path

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PPT = ROOT / "DPW_PPT.pptx"
BACKUP = ROOT / "DPW_PPT_before_speaker4_update.pptx"
OUT = ROOT / "DPW_PPT.pptx"
SHOT_DIR = ROOT / "data" / "ppt_real_screenshots"
CROP_DIR = SHOT_DIR / "crops"
BG_IMAGE = Path(
    r"E:\WeChatData\xwechat_files\wxid_xz3goysxc7an21_7c89\temp\RWTemp\2026-05\9ffabf8a760088ee6e67381d636eeba8\9fb1c1458dcee08893b557c20d589277.jpg"
)

BG = "07111F"
PANEL = "0F1C2D"
CARD = "122033"
BORDER = "27364D"
TEXT = "F2F6FF"
MUTED = "9AA9BD"
BLUE = "3B82F6"
TEAL = "2DD4BF"
PURPLE = "A855F7"
AMBER = "FBBF24"
GREEN = "86EFAC"
ROSE = "FB7185"


def rgb(h: str) -> RGBColor:
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def clear(slide) -> None:
    sp_tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        sp_tree.remove(shape._element)
    # Removing animated shapes without removing timing nodes can leave
    # PowerPoint animation references pointing at missing shape IDs.
    for child in list(slide._element):
        if child.tag.endswith("}timing"):
            slide._element.remove(child)


def add_rect(slide, x, y, w, h, fill=CARD, line=BORDER, transparency=0, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.fill.transparency = transparency
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(0.6)
    return shape


def add_text(slide, text, x, y, w, h, size=14, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    p.font.name = "Arial"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    p.text = text
    return box


def add_title(slide, section, title, subtitle, num):
    add_text(slide, section.upper(), 0.55, 0.30, 3.8, 0.22, 9, BLUE, True)
    add_text(slide, f"{num:02d} / 23", 11.65, 0.30, 1.0, 0.22, 9, MUTED, True, PP_ALIGN.RIGHT)
    add_text(slide, title, 0.55, 0.62, 9.2, 0.48, 25, TEXT, True)
    if subtitle:
        add_text(slide, subtitle, 0.57, 1.12, 10.9, 0.34, 11.5, MUTED)


def add_card(slide, x, y, w, h, label, value, note="", accent=BLUE, value_size=18):
    add_rect(slide, x, y, w, h, CARD, BORDER, 0, True)
    add_text(slide, label.upper(), x + 0.14, y + 0.12, w - 0.28, 0.18, 7.5, MUTED, True)
    add_text(slide, value, x + 0.14, y + 0.37, w - 0.28, 0.34, value_size, accent, True)
    if note:
        add_text(slide, note, x + 0.14, y + 0.76, w - 0.28, h - 0.78, 8.2, MUTED)


def add_footer(slide, role="Chi-Square Independence Tests"):
    add_rect(slide, 0.55, 6.86, 12.25, 0.36, PANEL, BORDER, 10, False)
    add_text(slide, "ZHOU Can · Speaker 4 · approx. 2 min", 0.72, 6.94, 3.1, 0.15, 8.2, TEAL, True)
    add_text(slide, role, 3.55, 6.94, 8.7, 0.15, 8.2, MUTED)


def add_bullets(slide, items, x, y, w, h, size=13, color=TEXT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(5)
    return box


def add_bg(slide, speaker4=False):
    if speaker4 and BG_IMAGE.exists():
        slide.shapes.add_picture(str(BG_IMAGE), Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
        overlay = add_rect(slide, 0, 0, 13.333, 7.5, BG, BG, 12, False)
        overlay.line.fill.background()
    else:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb(BG)


def add_picture(slide, path: Path, x, y, w, h, line=True):
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    if line:
        pic.line.color.rgb = rgb(BORDER)
        pic.line.width = Pt(0.8)
    return pic


def crop_to_ratio(src: Path, dst: Path, box, ratio=16 / 9):
    im = Image.open(src).convert("RGB")
    if box:
        im = im.crop(box)
    w, h = im.size
    current = w / h
    if current > ratio:
        new_w = int(h * ratio)
        left = (w - new_w) // 2
        im = im.crop((left, 0, left + new_w, h))
    elif current < ratio:
        new_h = int(w / ratio)
        top = max(0, (h - new_h) // 2)
        im = im.crop((0, top, w, top + new_h))
    dst.parent.mkdir(parents=True, exist_ok=True)
    im.save(dst, quality=96)
    return dst


def prepare_crops() -> dict[str, Path]:
    crops = {}
    main_box = (340, 285, 1740, 1073)
    for name in [
        "chi_genre_profitability",
        "chi_budget_profitability",
        "chi_release_profitability",
        "chi_genre_rating",
        "chi_budget_rating",
    ]:
        crops[name] = crop_to_ratio(
            SHOT_DIR / f"{name}.png",
            CROP_DIR / f"{name}_main.jpg",
            main_box,
        )
    crops["chi_overview"] = crop_to_ratio(
        SHOT_DIR / "chi_square_view.png",
        CROP_DIR / "chi_square_main.jpg",
        main_box,
    )
    crops["overview"] = crop_to_ratio(
        SHOT_DIR / "overview_full.png",
        CROP_DIR / "overview_view.jpg",
        (0, 0, 1760, 990),
    )
    crops["regression"] = crop_to_ratio(
        SHOT_DIR / "regression_view.png",
        CROP_DIR / "regression_view.jpg",
        (330, 230, 1740, 1025),
    )
    crops["export"] = crop_to_ratio(
        SHOT_DIR / "export_view.png",
        CROP_DIR / "export_view.jpg",
        (330, 230, 1740, 1025),
    )
    return crops


def format_p(p, size=9, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    p.font.name = "Arial"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(color)
    p.alignment = align


def add_native_table(slide, rows, x, y, w, h, widths=None, font_size=8.5):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if widths:
        for idx, width in enumerate(widths):
            table.columns[idx].width = Inches(width)
    for r, row in enumerate(rows):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(PANEL if r == 0 else (CARD if r % 2 else BG))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                format_p(p, size=font_size if r else font_size - 0.2, color=MUTED if r == 0 else TEXT, bold=(r == 0))
    return table_shape


def update_talk_map(slide):
    clear(slide)
    add_bg(slide)
    add_title(slide, "Talk Map", "Presentation Flow Aligned with Team Responsibilities",
              "A 15-minute technical narrative followed by a 5-minute live demo and Q&A.", 2)
    rows = [
        ["01", "Project & Data Cleaning", "2.5 min", "Problem framing, TMDB source data, SQL-to-CSV cleaning pipeline"],
        ["02", "Database & ER Design", "2.5 min", "Relational schema, core tables, many-to-many links, SQL export"],
        ["03", "Dashboard Exploration", "3 min", "Overview, Revenue, Genres and Time views with real filters"],
        ["04", "Chi-Square Tests", "2 min", "Categorical feature engineering, independence tests, residual diagnostics"],
        ["05", "Regression & Prediction", "2.5 min", "Linear regression, Gradient Boosting, comparable-film output"],
        ["06", "Insights, Export & Close", "2.5 min", "Automated insight cards, CSV export, limitations and future work"],
    ]
    add_native_table(slide, [["#", "Section", "Time", "Presentation focus"]] + rows,
                     0.75, 1.68, 11.8, 3.95, widths=[0.65, 2.95, 1.05, 7.15], font_size=9)
    add_card(slide, 0.9, 5.95, 3.1, 0.72, "Demo Route", "Filters → Charts → Prediction → Q&A",
             "Speaker 3 or Speaker 6 operates; module owners answer follow-up questions.", AMBER, 12)
    add_card(slide, 4.3, 5.95, 3.1, 0.72, "Structure Principle", "Data pipeline before results",
             "The story moves from reproducibility to interpretation.", TEAL, 12)
    add_card(slide, 7.7, 5.95, 3.1, 0.72, "Style", "Engineering presentation",
             "Fewer slogans, more evidence, consistent terminology.", PURPLE, 12)


def update_talk_map_in_place(slide):
    mapping = {
        "02 / 18": "02 / 23",
        "按演讲分工重构后的汇报路线": "Presentation Flow by Team Responsibility",
        "顺序跟随 dashboard 侧边栏视图，同时让每位成员有清晰负责段落。": "The sequence follows the dashboard sidebar and gives each speaker a focused technical segment.",
        "Project + Data Cleaning": "Project",
        "2.5 min\n背景、数据来源、Python 清洗": "2.5 min\nProblem, source data, cleaning",
        "Database + ER": "Database",
        "2.5 min\n关系设计、核心表、SQL 导出": "2.5 min\nER schema, core tables, SQL export",
        "Dashboard Views": "Dashboard",
        "3 min\nOverview / Revenue / Genres / Time": "3 min\nOverview, Revenue, Genres, Time",
        "2 min\n分类变量显著性检验": "2 min\nIndependence tests and residuals",
        "2.5 min\n线性回归与 Gradient Boosting": "2.5 min\nRegression and Gradient Boosting",
        "Insights + Export + Summary": "Insights",
        "2.5 min\n自动洞察、导出、局限与未来方向": "2.5 min\nFindings, export, limits, future work",
        "Demo & Q&A: 5 min · 第三位或第六位主操作，其他成员辅助回答。": "Demo & Q&A: 5 min · Speaker 3 or Speaker 6 operates the system; module owners support technical questions.",
    }
    replace_text(slide, mapping)
    stage_titles = {"Project", "Database", "Dashboard", "Chi-Square", "Regression", "Insights"}
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip()
            if text in stage_titles:
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(25)
                        run.font.name = "Arial"
            elif text.startswith("2.5 min") or text.startswith("3 min") or text.startswith("2 min"):
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(12)
                        run.font.name = "Arial"
    # Mask the original narrow cards before placing the cleaner compact map.
    # The source slide has fragile PowerPoint object references, so this keeps
    # existing objects in place and adds a stable editable overlay above them.
    bg_mask = add_rect(slide, 0.0, 1.58, 13.333, 5.18, BG, BG, 0, False)
    bg_mask.line.fill.background()
    compact = [
        ("01", "Project + Data Cleaning", "2.5 min", "Problem framing, TMDB source data, SQL-to-CSV cleaning.", BLUE),
        ("02", "Database + ER Design", "2.5 min", "Relational schema, core tables, SQL export path.", TEAL),
        ("03", "Dashboard Exploration", "3 min", "Overview, Revenue, Genres and Time views.", PURPLE),
        ("04", "Chi-Square Tests", "2 min", "Categorical hypotheses, p-values and residuals.", AMBER),
        ("05", "Regression + Prediction", "2.5 min", "Linear fit, model metrics and comparable-film output.", GREEN),
        ("06", "Insights + Export", "2.5 min", "Automated findings, CSV export, limits and future work.", ROSE),
    ]
    x_positions = [0.74, 4.76, 8.78]
    y_positions = [1.95, 4.12]
    card_w, card_h = 3.52, 1.34
    for i, (num, title, time, note, color) in enumerate(compact):
        x = x_positions[i % 3]
        y = y_positions[i // 3]
        add_rect(slide, x, y, card_w, card_h, CARD, BORDER, 0, False)
        add_text(slide, num, x + 0.18, y + 0.16, 0.44, 0.22, 9.5, MUTED, True)
        add_text(slide, title, x + 0.72, y + 0.13, card_w - 0.92, 0.28, 14.5, color, True)
        add_text(slide, time, x + 0.72, y + 0.49, 0.9, 0.2, 10.2, MUTED, True)
        add_text(slide, note, x + 0.72, y + 0.78, card_w - 0.98, 0.38, 8.9, MUTED)
    add_rect(slide, 0.74, 6.05, 11.56, 0.42, PANEL, BORDER, 10, False)
    add_text(
        slide,
        "Demo & Q&A: 5 min · Speaker 3 or Speaker 6 operates the system; module owners support technical questions.",
        0.94, 6.17, 11.1, 0.16, 8.8, TEXT, True, PP_ALIGN.CENTER,
    )


def update_speaker3_text_and_images(slide8, slide9, slide10, crops):
    replacements = {
        "SPEAKER 3": "Speaker 3",
        "07 / 18": "08 / 23",
        "08 / 18": "09 / 23",
        "09 / 18": "10 / 23",
        "系统整体界面：一个入口，八个分析视图": "System Interface: One Entry Point, Eight Analytical Views",
        "这一段按侧边栏顺序讲 Overview / Revenue / Genres / Time。": "This section follows the sidebar order: Overview, Revenue, Genres and Time.",
        "第三位 · 约 3 min": "Speaker 3 · approx. 3 min",
        "Revenue 视图：预算推高票房，但不保证 ROI": "Revenue View: Budget Raises Revenue, Not Necessarily ROI",
        "先用预算-票房回答商业规模，再用预算-ROI 解释回报率风险。": "Use budget-revenue for market scale, then budget-ROI for efficiency and risk.",
        "Genres 与 Time：类型和档期决定盈利侧重点": "Genres and Time: Profitability Depends on Category and Release Window",
        "把 dashboard 的类型排序、评分分布和上映月份曲线合并成一段讲。": "Combine genre ranking, rating distribution and release-month revenue into one analytical readout.",
    }
    for slide in [slide8, slide9, slide10]:
        replace_text(slide, replacements)
    # Replace the old GUI screenshot on slide 8 with a fresh running-system screenshot.
    for shape in list(slide8.shapes):
        if shape.shape_type == 13 and shape.left > Inches(0.5) and shape.top > Inches(1):
            slide8.shapes._spTree.remove(shape._element)
    add_picture(slide8, crops["overview"], 0.75, 1.45, 8.25, 4.65)


def replace_text(slide, mapping):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    text = r.text
                    for old, new in mapping.items():
                        text = text.replace(old, new)
                    r.text = text


def replace_slide_number(slide, idx, total):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            text = shape.text.strip()
            if re.fullmatch(r"\d{2}\s*/\s*\d+", text):
                shape.text = f"{idx:02d} / {total}"


def rebuild_slide_11(slide, crops):
    clear(slide)
    add_bg(slide, True)
    add_title(slide, "ZHOU Can", "Chi-Square Independence Testing Framework",
              "The module converts filtered movie records into categorical hypotheses and evaluates whether observed distributions deviate from independence.", 11)
    add_footer(slide, "Method overview")
    add_picture(slide, crops["chi_overview"], 6.55, 1.55, 5.85, 3.3)
    add_bullets(slide, [
        "Purpose: test whether two categorical variables are statistically associated under the current dashboard filters.",
        "Null hypothesis H0: the variables are independent; p < 0.05 rejects H0.",
        "Outputs: contingency table, chi-square statistic, degrees of freedom, p-value, Cramer's V and standardized residuals.",
    ], 0.72, 1.65, 5.45, 2.05, 13.4)
    steps = [
        ("01", "Filter data", "Year, genre and vote thresholds define the analysis slice."),
        ("02", "Derive categories", "budget_tier, season, profitable and rating_tier."),
        ("03", "Run test", "scipy.stats.chi2_contingency over observed counts."),
        ("04", "Diagnose cells", "Residual heatmap explains which cells drive deviation."),
    ]
    for i, (n, label, note) in enumerate(steps):
        x = 0.75 + i * 3.05
        add_card(slide, x, 5.15, 2.72, 0.92, n, label, note, [BLUE, TEAL, PURPLE, AMBER][i], 12)


def rebuild_slide_12(slide):
    clear(slide)
    add_bg(slide, True)
    add_title(slide, "ZHOU Can", "Derived Variables and Test Configuration",
              "Continuous movie attributes are discretized into interpretable categories before entering the chi-square tests.", 12)
    add_footer(slide, "Feature engineering for categorical inference")
    rows = [
        ["Derived field", "Definition", "Why it matters"],
        ["profitable", "revenue > budget", "Turns financial outcome into a binary response variable."],
        ["budget_tier", "<$10M, $10-50M, $50-150M, >$150M", "Makes production scale comparable across films."],
        ["season", "Winter, Spring, Summer, Fall, Awards Season", "Maps release month to market timing strategy."],
        ["rating_tier", "Poor, Average, Good, Great", "Converts a continuous score into audience-quality bands."],
    ]
    add_native_table(slide, rows, 0.75, 1.65, 7.15, 3.4, widths=[1.55, 2.55, 3.05], font_size=8.7)
    add_card(slide, 8.35, 1.68, 3.0, 0.95, "Testing Engine", "chi2_contingency", "SciPy implementation over contingency tables.", BLUE, 14)
    add_card(slide, 8.35, 2.95, 3.0, 0.95, "Effect Size", "Cramer's V", "Normalizes the chi-square statistic for association strength.", TEAL, 14)
    add_card(slide, 8.35, 4.22, 3.0, 0.95, "Diagnostics", "Residual heatmap", "Blue cells occur more often than expected; red cells occur less often.", PURPLE, 14)
    add_bullets(slide, [
        "The same computation runs after every filter change.",
        "This keeps the statistics consistent with the visible dashboard slice.",
    ], 0.95, 5.55, 10.6, 0.75, 13)


def rebuild_slide_13(slide, crops):
    clear(slide)
    add_bg(slide, True)
    add_title(slide, "ZHOU Can", "Profitability Hypotheses: No Significant Association in the Default Slice",
              "Genre, budget tier and release season do not show statistically significant association with profitability under the default filters.", 13)
    add_footer(slide, "Profitability tests")
    add_picture(slide, crops["chi_release_profitability"], 0.72, 1.55, 6.95, 3.9)
    rows = [
        ["Test", "n", "chi-square", "p-value", "Cramer's V", "Decision"],
        ["Genre x Profitability", "1,026", "4.34", "0.7402", "0.065", "Fail to reject H0"],
        ["Budget Tier x Profitability", "1,026", "1.13", "0.7703", "0.033", "Fail to reject H0"],
        ["Release Season x Profitability", "1,026", "2.42", "0.6589", "0.049", "Fail to reject H0"],
    ]
    add_native_table(slide, rows, 7.95, 1.62, 4.45, 2.55,
                     widths=[1.38, 0.48, 0.72, 0.72, 0.72, 1.15], font_size=7.2)
    add_bullets(slide, [
        "Interpretation: the observed profitability rates are close to the expected counts under independence.",
        "Engineering note: this conclusion is conditional on the active filters; changing genres or vote thresholds recomputes the test.",
    ], 8.05, 4.52, 4.25, 1.15, 11.2)


def rebuild_slide_14(slide, crops):
    clear(slide)
    add_bg(slide, True)
    add_title(slide, "ZHOU Can", "Rating-Tier Hypotheses: Strong Evidence of Association",
              "Rating distributions vary meaningfully across genres and budget tiers, with small-to-moderate effect sizes.", 14)
    add_footer(slide, "Rating-tier tests")
    add_picture(slide, crops["chi_budget_rating"], 0.72, 1.55, 6.95, 3.9)
    rows = [
        ["Test", "n", "chi-square", "p-value", "Cramer's V", "Decision"],
        ["Genre x Rating Tier", "956", "106.16", "<0.0001", "0.192", "Reject H0"],
        ["Budget Tier x Rating Tier", "946", "110.30", "<0.0001", "0.197", "Reject H0"],
    ]
    add_native_table(slide, rows, 7.95, 1.62, 4.45, 1.9,
                     widths=[1.5, 0.48, 0.72, 0.72, 0.72, 1.0], font_size=7.5)
    add_card(slide, 8.0, 3.86, 1.9, 0.86, "Significance", "p < 0.001", "extremely strong evidence", GREEN, 13)
    add_card(slide, 10.15, 3.86, 1.9, 0.86, "Effect size", "~0.20", "small-to-moderate association", TEAL, 13)
    add_bullets(slide, [
        "The rating module is not just descriptive: it identifies systematic category-level distribution shifts.",
        "Residuals show which budget tiers are over- or under-represented in high-rating bands.",
    ], 8.02, 5.0, 4.3, 0.95, 10.8)


def rebuild_slide_15(slide, crops):
    clear(slide)
    add_bg(slide, True)
    add_title(slide, "ZHOU Can", "Residual Diagnostics: Explaining Where the Association Comes From",
              "The heatmap translates the chi-square result into interpretable cell-level deviations from expected counts.", 15)
    add_footer(slide, "Residual heatmap and distribution view")
    add_picture(slide, crops["chi_genre_rating"], 0.72, 1.55, 6.95, 3.9)
    add_card(slide, 8.0, 1.65, 3.1, 0.9, "Standardized residual", "(observed - expected) / sqrt(expected)", "cell-level diagnostic", BLUE, 11)
    add_card(slide, 8.0, 2.85, 3.1, 0.9, "Rule of thumb", "|residual| > 2", "meaningful over/under-representation", AMBER, 13)
    add_card(slide, 8.0, 4.05, 3.1, 0.9, "Readout", "Blue vs red cells", "above vs below independence expectation", PURPLE, 13)
    add_bullets(slide, [
        "The stacked proportion chart checks whether category distributions are visually similar or structurally different.",
        "Together, the KPI row and residual plot prevent overclaiming from p-value alone.",
    ], 8.05, 5.25, 4.25, 0.85, 10.8)


def rebuild_slide_16(slide):
    clear(slide)
    add_bg(slide, True)
    add_title(slide, "ZHOU Can", "Engineering Implementation and Takeaway",
              "The chi-square view is implemented as a reusable statistical workflow inside the Streamlit dashboard.", 16)
    add_footer(slide, "Implementation summary")
    add_bullets(slide, [
        "`_add_derived_cols(df)` constructs the categorical variables without mutating the original DataFrame.",
        "`_run_chi2(df, row_col, col_col)` builds a contingency table, runs SciPy's chi-square test and returns standardized residuals.",
        "`_interpret_p(p)` maps the p-value to an explanation that can be displayed directly in the dashboard.",
        "Plotly heatmaps and stacked bars provide both statistical significance and visual diagnostics.",
    ], 0.82, 1.62, 6.0, 2.95, 14)
    add_card(slide, 7.25, 1.7, 2.25, 1.0, "Reusable", "5 tests", "same computation pattern", BLUE, 16)
    add_card(slide, 9.85, 1.7, 2.25, 1.0, "Interactive", "filter-aware", "results recompute live", TEAL, 16)
    add_card(slide, 7.25, 3.05, 2.25, 1.0, "Transparent", "counts + p-value", "not a black-box result", PURPLE, 16)
    add_card(slide, 9.85, 3.05, 2.25, 1.0, "Takeaway", "ratings vary", "profitability less clear", GREEN, 16)
    add_text(slide, "Final message: the module separates weak profitability evidence from strong rating-distribution evidence, which is exactly why formal testing is useful in the dashboard.", 1.0, 5.65, 11.1, 0.42, 16, TEXT, True, PP_ALIGN.CENTER)


def replace_ui_screenshots(prs, crops):
    # Slide 17 regression screenshot.
    if len(prs.slides) >= 17:
        slide = prs.slides[16]
        for shape in list(slide.shapes):
            if shape.shape_type == 13 and shape.top > Inches(1):
                slide.shapes._spTree.remove(shape._element)
        add_picture(slide, crops["regression"], 0.75, 1.45, 7.2, 3.8)
    # Slide 21 export screenshot.
    if len(prs.slides) >= 21:
        slide = prs.slides[20]
        for shape in list(slide.shapes):
            if shape.shape_type == 13 and shape.top > Inches(1) and shape.left > Inches(0.5):
                slide.shapes._spTree.remove(shape._element)
        add_picture(slide, crops["export"], 1.14, 1.35, 11.08, 5.53)


def global_text_cleanup(prs):
    total = len(prs.slides)
    global_map = {
        "：": ":",
        "SPEAKER 4": "ZHOU Can",
        "speaker 4": "ZHOU Can",
        "第四位": "ZHOU Can",
        "约": "approx.",
        "系统整体界面: 一个入口，八个分析视图": "System Interface: One Entry Point, Eight Analytical Views",
        "这一段按侧边栏顺序讲 Overview / Revenue / Genres / Time。": "This section follows the sidebar order: Overview, Revenue, Genres and Time.",
        "Revenue 视图: 预算推高票房，但不保证 ROI": "Revenue View: Budget Raises Revenue, Not Necessarily ROI",
        "先用预算-票房回答商业规模，再用预算-ROI 解释回报率风险。": "Use budget-revenue for market scale, then budget-ROI for efficiency and risk.",
        "Genres 与 Time: 类型和档期决定盈利侧重点": "Genres and Time: Profitability Depends on Category and Release Window",
        "把 dashboard 的类型排序、评分分布和上映月份曲线合并成一段讲。": "Combine genre ranking, rating distribution and release-month revenue into one analytical readout.",
        "派生字段与解释方式": "Derived variables and interpretation",
        "讲法:先说明所有筛选会同步刷新 KPI、图表和导出结果。": "Presenter cue: explain that every filter updates KPIs, charts and exports together.",
    }
    for idx, slide in enumerate(prs.slides, 1):
        replace_text(slide, global_map)
        replace_slide_number(slide, idx, total)


def check_chinese(prs) -> list[tuple[int, str]]:
    hits = []
    pattern = re.compile(r"[\u4e00-\u9fff]")
    for idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                text = "\n".join(p.text for p in shape.text_frame.paragraphs)
                if pattern.search(text):
                    hits.append((idx, text))
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        if pattern.search(cell.text):
                            hits.append((idx, cell.text))
    return hits


def main():
    if not BACKUP.exists():
        shutil.copy2(PPT, BACKUP)
    crops = prepare_crops()
    prs = Presentation(PPT)
    update_talk_map_in_place(prs.slides[1])
    update_speaker3_text_and_images(prs.slides[7], prs.slides[8], prs.slides[9], crops)
    rebuild_slide_11(prs.slides[10], crops)
    rebuild_slide_12(prs.slides[11])
    rebuild_slide_13(prs.slides[12], crops)
    rebuild_slide_14(prs.slides[13], crops)
    rebuild_slide_15(prs.slides[14], crops)
    rebuild_slide_16(prs.slides[15])
    replace_ui_screenshots(prs, crops)
    global_text_cleanup(prs)
    hits = check_chinese(prs)
    if hits:
        for idx, text in hits:
            print(f"Chinese text remains on slide {idx}: {text[:120]}")
        raise SystemExit("Chinese text remains")
    prs.save(OUT)
    print(f"updated {OUT}")
    print(f"backup {BACKUP}")


if __name__ == "__main__":
    main()
