from __future__ import annotations

import math
import random
import sys
from pathlib import Path

import numpy as np
import pandas as pd
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))
OUT = ROOT / "DPW_PPT.pptx"
FALLBACK_OUT = ROOT / "DPW_PPT_revised.pptx"
ASSET_DIR = ROOT / "data" / "ppt_revision_assets"
CSV = ROOT / "data" / "processed_movies.csv"
METRIC_MD = ROOT / "DPW_PPT_metric_check.md"

BG = "07111F"
PANEL = "0F1C2D"
CARD = "122033"
BORDER = "27364D"
TEXT = "F2F6FF"
MUTED = "9AA9BD"
BLUE = "3B82F6"
TEAL = "2DD4BF"
PURPLE = "A855F7"
AMBER = "FBBF24"
GREEN = "86EFAC"
ROSE = "FB7185"

MONTHS = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
CORE_GENRES = [
    "Action", "Adventure", "Animation", "Comedy", "Crime", "Drama", "Family",
    "Fantasy", "Horror", "Mystery", "Romance", "Science Fiction", "Thriller",
]


def rgb(hex_value: str) -> RGBColor:
    h = hex_value.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def money(value: float) -> str:
    if pd.isna(value):
        return "$0"
    if abs(value) >= 1_000_000_000:
        return f"${value / 1_000_000_000:.2f}B"
    if abs(value) >= 1_000_000:
        return f"${value / 1_000_000:.1f}M"
    if abs(value) >= 1_000:
        return f"${value / 1_000:.1f}K"
    return f"${value:.0f}"


def pct(value: float) -> str:
    return f"{value * 100:.1f}%"


def load_df() -> pd.DataFrame:
    df = pd.read_csv(CSV)
    for col in ["budget", "revenue", "year", "release_month", "runtime", "avg_rating", "roi", "analysis_vote_count"]:
        df[col] = pd.to_numeric(df[col], errors="coerce")
    df["genres"] = df["genres"].fillna("")
    df["primary_genre"] = df["primary_genre"].fillna("Unknown")
    return df


def explode_genres(df: pd.DataFrame) -> pd.DataFrame:
    genre_df = df.copy()
    genre_df["genre"] = genre_df["genres"].str.split("|")
    genre_df = genre_df.explode("genre")
    genre_df["genre"] = genre_df["genre"].replace("", pd.NA)
    return genre_df.dropna(subset=["genre"])


def compute_metrics(df: pd.DataFrame) -> dict:
    valid = df[(df["budget"] > 0) & (df["revenue"] > 0) & df["roi"].notna()].copy()
    rated = df[df["avg_rating"].notna()].copy()
    genre_df = explode_genres(valid)
    genre_df = genre_df[genre_df["genre"].isin(CORE_GENRES)]
    genre_roi = genre_df.groupby("genre")["roi"].median().sort_values(ascending=False)
    genre_rating_df = explode_genres(rated)
    genre_rating_df = genre_rating_df[genre_rating_df["genre"].isin(CORE_GENRES)]
    genre_rating = genre_rating_df.groupby("genre")["avg_rating"].median().sort_values(ascending=False)
    month_revenue = valid.groupby("release_month")["revenue"].mean().sort_values(ascending=False)
    roi_cap = valid["roi"].quantile(0.98)
    roi_corr_df = valid[(valid["roi"] <= roi_cap) & (valid["budget"] > 0)]
    profitable = valid[valid["revenue"] > valid["budget"]]
    unprofitable = valid[valid["revenue"] <= valid["budget"]]
    low_roi = valid[valid["budget"] < 10_000_000]["roi"].median()
    blockbuster_roi = valid[valid["budget"] > 150_000_000]["roi"].median()
    rating_bins = pd.cut(valid["avg_rating"], bins=[0, 5, 6, 7, 8, 10], labels=["<5", "5-6", "6-7", "7-8", "8+"])
    rating_rev = valid.groupby(rating_bins, observed=False)["revenue"].mean()
    multi_genre_share = (df["genres"].str.contains(r"\|", regex=True).sum() / len(df))

    return {
        "df": df,
        "valid": valid,
        "rated": rated,
        "movie_count": len(df),
        "valid_roi_count": len(valid),
        "rated_count": len(rated),
        "budget_revenue_corr": valid[["budget", "revenue"]].corr().iloc[0, 1],
        "budget_roi_corr": roi_corr_df[["budget", "roi"]].corr().iloc[0, 1],
        "top_genre": genre_roi.index[0],
        "top_genre_roi": genre_roi.iloc[0],
        "bottom_genre": genre_roi.index[-1],
        "bottom_genre_roi": genre_roi.iloc[-1],
        "genre_roi": genre_roi,
        "best_month": int(month_revenue.index[0]),
        "best_month_rev": month_revenue.iloc[0],
        "worst_month": int(month_revenue.index[-1]),
        "worst_month_rev": month_revenue.iloc[-1],
        "month_revenue": month_revenue.sort_index(),
        "low_roi": low_roi,
        "blockbuster_roi": blockbuster_roi,
        "rating_top_genre": genre_rating.index[0],
        "rating_top": genre_rating.iloc[0],
        "rating_bottom_genre": genre_rating.index[-1],
        "rating_bottom": genre_rating.iloc[-1],
        "profit_rate": len(profitable) / len(valid),
        "profit_rating_diff": profitable["avg_rating"].mean() - unprofitable["avg_rating"].mean(),
        "long_roi": valid[valid["runtime"] > 150]["roi"].median(),
        "rating_sweet_spot": str(rating_rev.idxmax()),
        "rating_sweet_spot_rev": rating_rev.max(),
        "multi_genre_share": multi_genre_share,
    }


def font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = [
        "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/segoeuib.ttf" if bold else "C:/Windows/Fonts/segoeui.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def draw_axes(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int]) -> None:
    x0, y0, x1, y1 = box
    grid = (154, 169, 189, 45)
    axis = (154, 169, 189, 100)
    for i in range(5):
        y = y0 + i * (y1 - y0) / 4
        draw.line((x0, y, x1, y), fill=grid, width=1)
    draw.line((x0, y1, x1, y1), fill=axis, width=2)
    draw.line((x0, y0, x0, y1), fill=axis, width=2)


def scale(values: pd.Series, vmin: float, vmax: float, out_min: int, out_max: int, log: bool = False) -> list[float]:
    vals = values.astype(float).clip(lower=max(vmin, 1) if log else vmin)
    if log:
        vals = np.log10(vals)
        vmin, vmax = math.log10(max(vmin, 1)), math.log10(max(vmax, 1))
    span = max(vmax - vmin, 1e-9)
    return (out_min + (vals - vmin) / span * (out_max - out_min)).tolist()


def chart_budget_revenue(m: dict) -> Path:
    valid = m["valid"].copy()
    sample = valid.sample(min(900, len(valid)), random_state=42)
    path = ASSET_DIR / "chart_budget_revenue.png"
    im = Image.new("RGB", (1500, 760), "#" + BG)
    d = ImageDraw.Draw(im, "RGBA")
    box = (105, 70, 1410, 650)
    draw_axes(d, box)
    x = scale(sample["budget"], 1, valid["budget"].quantile(0.99), box[0], box[2], log=True)
    y = scale(sample["revenue"], 1, valid["revenue"].quantile(0.99), box[3], box[1], log=True)
    for px, py in zip(x, y):
        d.ellipse((px - 3, py - 3, px + 3, py + 3), fill=(59, 130, 246, 145))
    xs = np.log10(valid["budget"].clip(lower=1))
    ys = np.log10(valid["revenue"].clip(lower=1))
    coef = np.polyfit(xs, ys, 1)
    line_x = np.array([xs.quantile(0.02), xs.quantile(0.99)])
    line_y = coef[0] * line_x + coef[1]
    lx = (box[0] + (line_x - math.log10(1)) / (math.log10(valid["budget"].quantile(0.99)) - math.log10(1)) * (box[2] - box[0]))
    ly = (box[3] + (line_y - math.log10(1)) / (math.log10(valid["revenue"].quantile(0.99)) - math.log10(1)) * (box[1] - box[3]))
    d.line((lx[0], ly[0], lx[1], ly[1]), fill=(242, 246, 255, 230), width=5)
    d.text((105, 22), f"Budget vs Revenue · Pearson r = {m['budget_revenue_corr']:.2f}", fill="#" + TEXT, font=font(34, True))
    d.text((105, 675), "Log-scaled axes · all films with budget > 0 and revenue > 0", fill="#" + MUTED, font=font(22))
    im.save(path)
    return path


def chart_genre_roi(m: dict) -> Path:
    s = m["genre_roi"].head(10).sort_values()
    path = ASSET_DIR / "chart_genre_roi.png"
    im = Image.new("RGB", (1500, 760), "#" + BG)
    d = ImageDraw.Draw(im, "RGBA")
    x0, y0, x1, y1 = 260, 90, 1370, 660
    maxv = s.max() * 1.15
    d.text((90, 25), "Median ROI by Genre · all valid ROI films", fill="#" + TEXT, font=font(34, True))
    row_h = (y1 - y0) / len(s)
    for i, (genre, value) in enumerate(s.items()):
        y = y0 + i * row_h + 8
        w = int((x1 - x0) * value / maxv)
        color = (59 + i * 8, 130, 246 + min(i * 2, 9), 235)
        d.text((70, y + 8), genre, fill="#" + TEXT, font=font(23, True))
        d.rounded_rectangle((x0, y, x0 + w, y + row_h - 16), radius=12, fill=color)
        d.text((x0 + w + 14, y + 8), f"{value:.2f}x", fill="#" + TEAL, font=font(23, True))
    im.save(path)
    return path


def chart_month_revenue(m: dict) -> Path:
    s = m["month_revenue"]
    path = ASSET_DIR / "chart_month_revenue.png"
    im = Image.new("RGB", (1500, 760), "#" + BG)
    d = ImageDraw.Draw(im, "RGBA")
    box = (100, 80, 1400, 640)
    draw_axes(d, box)
    vals = [s.get(i, 0) for i in range(1, 13)]
    ymax = max(vals) * 1.15
    pts = []
    for i, val in enumerate(vals):
        x = box[0] + i * (box[2] - box[0]) / 11
        y = box[3] - val / ymax * (box[3] - box[1])
        pts.append((x, y))
    area = [(box[0], box[3])] + pts + [(box[2], box[3])]
    d.polygon(area, fill=(45, 212, 191, 42))
    d.line(pts, fill=(45, 212, 191, 245), width=6, joint="curve")
    for i, (x, y) in enumerate(pts):
        d.ellipse((x - 8, y - 8, x + 8, y + 8), fill=(110, 231, 183, 245))
        d.text((x - 18, box[3] + 18), MONTHS[i], fill="#" + MUTED, font=font(20))
    d.text((100, 25), f"Average Revenue by Release Month · peak: {MONTHS[m['best_month']-1]}", fill="#" + TEXT, font=font(34, True))
    im.save(path)
    return path


def chart_budget_roi(m: dict) -> Path:
    valid = m["valid"].copy()
    cap = valid["roi"].quantile(0.98)
    plot = valid[valid["roi"] <= cap].sample(min(900, len(valid)), random_state=7)
    path = ASSET_DIR / "chart_budget_roi.png"
    im = Image.new("RGB", (1500, 760), "#" + BG)
    d = ImageDraw.Draw(im, "RGBA")
    box = (105, 70, 1410, 650)
    draw_axes(d, box)
    x = scale(plot["budget"], 1, valid["budget"].quantile(0.99), box[0], box[2], log=True)
    y = scale(plot["roi"], 0, cap, box[3], box[1])
    for px, py in zip(x, y):
        d.ellipse((px - 3, py - 3, px + 3, py + 3), fill=(163, 230, 53, 135))
    d.text((105, 22), f"Budget vs ROI · low-budget median {m['low_roi']:.2f}x vs blockbuster {m['blockbuster_roi']:.2f}x", fill="#" + TEXT, font=font(31, True))
    d.text((105, 675), "Budget is log-scaled · ROI capped at 98th percentile for readability", fill="#" + MUTED, font=font(22))
    im.save(path)
    return path


def chart_rating_distribution(m: dict) -> Path:
    rated = explode_genres(m["rated"])
    rated = rated[rated["genre"].isin(CORE_GENRES)]
    order = rated.groupby("genre")["avg_rating"].median().sort_values(ascending=False).head(11).index.tolist()
    path = ASSET_DIR / "chart_rating_distribution.png"
    im = Image.new("RGB", (1500, 760), "#" + BG)
    d = ImageDraw.Draw(im, "RGBA")
    box = (90, 90, 1410, 620)
    draw_axes(d, box)
    for i, genre in enumerate(order):
        vals = rated.loc[rated["genre"] == genre, "avg_rating"].dropna()
        q1, med, q3 = vals.quantile([0.25, 0.5, 0.75])
        lo, hi = vals.quantile([0.05, 0.95])
        x = box[0] + (i + 0.5) * (box[2] - box[0]) / len(order)
        def sy(v): return box[3] - v / 10 * (box[3] - box[1])
        d.line((x, sy(lo), x, sy(hi)), fill=(154, 169, 189, 160), width=3)
        d.rounded_rectangle((x - 34, sy(q3), x + 34, sy(q1)), radius=8, fill=(168, 85, 247, 165), outline=(242, 246, 255, 120))
        d.line((x - 38, sy(med), x + 38, sy(med)), fill=(251, 191, 36, 235), width=4)
        label = "Sci-Fi" if genre == "Science Fiction" else genre
        d.text((x - 38, box[3] + 18), label[:10], fill="#" + MUTED, font=font(17))
    d.text((90, 25), "Rating Distribution by Genre · 0-10 scale", fill="#" + TEXT, font=font(34, True))
    im.save(path)
    return path


def write_metric_check(m: dict) -> None:
    lines = [
        "# DPW_PPT metric check",
        "",
        "All headline metrics use the current `data/processed_movies.csv`.",
        "",
        "| Metric | Value | Deck usage |",
        "|---|---:|---|",
        f"| Movies written | {m['movie_count']:,} | Dataset slide |",
        f"| Valid ROI films | {m['valid_roi_count']:,} | Dataset and chart footers |",
        f"| Films with ratings | {m['rated_count']:,} | Dataset/rating slides |",
        f"| Budget-Revenue Pearson r | {m['budget_revenue_corr']:.2f} | Budget vs revenue |",
        f"| Budget-ROI Pearson r, 98th percentile cap | {m['budget_roi_corr']:.2f} | Budget vs ROI |",
        f"| Top median ROI genre | {m['top_genre']} {m['top_genre_roi']:.2f}x | Genre slide |",
        f"| Bottom median ROI genre | {m['bottom_genre']} {m['bottom_genre_roi']:.2f}x | Genre slide |",
        f"| Best average revenue month | {MONTHS[m['best_month'] - 1]} {money(m['best_month_rev'])} | Timing slide |",
        f"| Worst average revenue month | {MONTHS[m['worst_month'] - 1]} {money(m['worst_month_rev'])} | Timing slide |",
        f"| Profitable film share | {pct(m['profit_rate'])} | What we learned |",
        f"| Profitable film rating lift | {m['profit_rating_diff']:.2f} points | What we learned |",
        f"| Long film median ROI | {m['long_roi']:.2f}x | What we learned |",
        f"| Highest average revenue rating band | {m['rating_sweet_spot']} {money(m['rating_sweet_spot_rev'])} | What we learned |",
    ]
    METRIC_MD.write_text("\n".join(lines) + "\n", encoding="utf-8")


def extract_old_media() -> dict[str, Path]:
    from zipfile import ZipFile

    source = ROOT / "DPW_PPT_before_revision.pptx"
    media = {}
    if not source.exists():
        source = OUT
    with ZipFile(source) as z:
        for name in z.namelist():
            if name.startswith("ppt/media/image") and name.endswith(".png"):
                out = ASSET_DIR / Path(name).name
                out.write_bytes(z.read(name))
                media[Path(name).stem] = out
    return media


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(BG)


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int = 20, color: str = TEXT,
             bold: bool = False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def add_title(slide, kicker: str, title: str, subtitle: str | None = None, num: int | None = None) -> None:
    add_text(slide, kicker.upper(), 0.55, 0.35, 2.7, 0.22, 9, BLUE, True)
    if num is not None:
        add_text(slide, f"{num:02d} / 16", 11.6, 0.35, 1.1, 0.22, 9, MUTED, True, PP_ALIGN.RIGHT)
    add_text(slide, title, 0.55, 0.66, 7.8, 0.55, 28, TEXT, True)
    if subtitle:
        add_text(slide, subtitle, 0.58, 1.17, 10.5, 0.34, 12, MUTED)


def add_card(slide, x: float, y: float, w: float, h: float, label: str, value: str, note: str = "", accent: str = BLUE):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(CARD)
    shape.line.color.rgb = rgb(BORDER)
    shape.line.width = Pt(0.8)
    add_text(slide, label.upper(), x + 0.18, y + 0.16, w - 0.36, 0.22, 8, MUTED, True)
    add_text(slide, value, x + 0.18, y + 0.44, w - 0.36, 0.44, 22, accent, True)
    if note:
        add_text(slide, note, x + 0.18, y + 0.93, w - 0.36, 0.34, 9, MUTED)


def add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, size: int = 16, color: str = TEXT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(8)


def add_image_fit(slide, path: Path, x: float, y: float, w: float, h: float) -> None:
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def build_deck(m: dict, charts: dict[str, Path], media: dict[str, Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def slide():
        s = prs.slides.add_slide(blank)
        set_bg(s)
        return s

    s = slide()
    add_text(s, "26DPWPROJECT · COURSE FINAL PRESENTATION", 0.75, 0.55, 7.5, 0.24, 10, BLUE, True)
    add_text(s, "TMDB Movie Analytics\nDashboard", 0.75, 1.25, 7.4, 1.65, 39, TEXT, True)
    add_text(s, "Budget predicts revenue, but timing and genre shape profitability.", 0.78, 3.12, 8.2, 0.35, 17, TEAL, True)
    add_text(s, "Built from 45,346 films in the TMDB SQL export, with Streamlit analytics and a GradientBoosting box-office predictor.", 0.78, 3.65, 8.5, 0.6, 15, MUTED)
    add_card(s, 0.78, 5.05, 2.05, 1.08, "Stack", "Python · SQL", "Streamlit · Plotly · sklearn", BLUE)
    add_card(s, 3.05, 5.05, 2.1, 1.08, "Dataset", "45,346", "analysis-ready movies", TEAL)
    add_card(s, 5.38, 5.05, 2.05, 1.08, "Team", "Group 26", "six-person build", PURPLE)
    add_text(s, "ZHOU Can · CHEN Ziming · DAI Ling · LIAO Shuaiyu · PAN Meihao · ZHU Junyu", 0.78, 6.55, 9.0, 0.3, 11, MUTED)

    s = slide()
    add_title(s, "Overview", "Agenda", "Four modules for a 15-minute final presentation.", 2)
    for i, (label, value, note, color) in enumerate([
        ("01", "Question", "What drives revenue, ROI, timing and ratings?", BLUE),
        ("02", "Data", "SQL dump to analysis-ready CSV and ER model.", TEAL),
        ("03", "Dashboard", "Five core charts answer the research questions.", PURPLE),
        ("04", "Model", "Prediction module and final findings.", AMBER),
    ]):
        x = 0.85 + i * 3.05
        add_card(s, x, 2.15, 2.55, 2.45, label, value, note, color)
    add_text(s, "Goal: show a reproducible database-to-dashboard pipeline, then use it to answer the SDS research questions.", 1.0, 5.45, 11.4, 0.45, 18, TEXT, True, PP_ALIGN.CENTER)

    s = slide()
    add_title(s, "Problem", "Research Questions", "The deck is organized around the SDS questions, not around implementation inventory.", 3)
    add_bullets(s, [
        "Q1 · Box office: how strongly does budget predict revenue?",
        "Q2 · Genres: which genres deliver stronger median ROI?",
        "Q3 · Timing: which release months produce higher average revenue?",
        "Q4 · Ratings: how do genre and rating bands relate to audience response and revenue?",
    ], 0.85, 1.85, 7.6, 3.0, 19)
    add_card(s, 9.2, 1.9, 2.9, 1.1, "Evidence", "5 charts", "one chart per question", TEAL)
    add_card(s, 9.2, 3.25, 2.9, 1.1, "Method", "SQL → CSV", "reproducible local parser", BLUE)
    add_card(s, 9.2, 4.6, 2.9, 1.1, "Output", "Dashboard", "interactive filters + exports", PURPLE)

    s = slide()
    add_title(s, "Data", "Dataset and Cleaning", "A TMDB SQL export becomes one analysis-ready movie table.", 4)
    add_card(s, 0.75, 1.8, 2.15, 1.2, "Movies loaded", "45,433", "raw rows in movie.sql", BLUE)
    add_card(s, 3.15, 1.8, 2.15, 1.2, "Movies written", f"{m['movie_count']:,}", "99.8% retained", TEAL)
    add_card(s, 5.55, 1.8, 2.15, 1.2, "Full ROI", f"{m['valid_roi_count']:,}", "budget + revenue > 0", PURPLE)
    add_card(s, 7.95, 1.8, 2.15, 1.2, "With ratings", f"{m['rated_count']:,}", "0.5-5 scaled to 0-10", AMBER)
    add_bullets(s, [
        "Parsed 23 SQL dump files directly; no MySQL server required.",
        "Dropped 87 movies without parseable release dates.",
        "Derived year, month, ROI, average rating, primary genre and vote counts.",
        "Dashboard reads `data/processed_movies.csv` directly for reproducible submission.",
    ], 0.9, 3.65, 10.9, 2.2, 16)

    s = slide()
    add_title(s, "Data", "Database Schema", "Four core tables power the analytics; the full schema keeps wider movie metadata available.", 5)
    if "image1" in media:
        add_image_fit(s, media["image1"], 0.9, 1.45, 11.5, 3.8)
    add_card(s, 0.9, 5.7, 2.5, 0.8, "Core tables", "movie · genres", "link_genres · rate", BLUE)
    add_card(s, 3.65, 5.7, 2.6, 0.8, "Full schema", "23 tables", "credits, cast, companies", TEAL)
    add_card(s, 6.5, 5.7, 2.9, 0.8, "Analytics scope", "4 tables", "focused for live dashboard", PURPLE)

    s = slide()
    add_title(s, "Engineering", "Pipeline and Stack", "SQL dumps in, Streamlit dashboard out; lightweight enough for one-click local launch.", 6)
    add_bullets(s, [
        "01 Source: `tmdb.sql/` with movie, genres, link_genres and rate tables.",
        "02 Parse: `scripts/build_dataset.py` reads INSERT statements directly.",
        "03 Derive: year, release_month, ROI, average rating and primary genre.",
        "04 Serve: `streamlit_app.py` loads CSV and renders Plotly charts.",
    ], 0.75, 1.55, 6.2, 3.3, 16)
    add_card(s, 7.35, 1.6, 2.1, 1.0, "UI", "Streamlit", "single-page app", BLUE)
    add_card(s, 9.75, 1.6, 2.1, 1.0, "Charts", "Plotly", "shared theme", TEAL)
    add_card(s, 7.35, 3.0, 2.1, 1.0, "Stats", "SciPy", "chi-square tests", PURPLE)
    add_card(s, 9.75, 3.0, 2.1, 1.0, "Models", "sklearn", "GradientBoosting", AMBER)
    add_text(s, "One-click launchers: `run_windows.bat` and `run_unix.sh` create venvs, install deps and start Streamlit.", 0.85, 5.85, 11.4, 0.45, 14, MUTED)

    s = slide()
    add_title(s, "Product", "Dashboard Overview", "A dark single-page analytics dashboard with shared filters, KPIs and five core charts.", 7)
    if "image2" in media:
        add_image_fit(s, media["image2"], 0.72, 1.35, 11.9, 5.35)
    add_text(s, "All filters update KPI cards, charts and exports together.", 0.8, 6.9, 7.8, 0.25, 11, MUTED)

    chart_slides = [
        ("Chart 1 of 5", "Budget predicts revenue", "Higher budgets generally raise box-office revenue, but the scatter still leaves room for both hits and expensive flops.", charts["budget_revenue"], [
            ("Pearson r", f"{m['budget_revenue_corr']:.2f}", f"n = {m['valid_roi_count']:,}", BLUE),
            ("Scope", "All valid ROI", "budget > 0, revenue > 0", TEAL),
        ]),
        ("Chart 2 of 5", "Genre profitability changes the story", "Median ROI avoids mega-hit distortion and shows genre profitability, not just absolute revenue.", charts["genre_roi"], [
            ("Top genre", m["top_genre"], f"{m['top_genre_roi']:.2f}x median ROI", PURPLE),
            ("Bottom genre", m["bottom_genre"], f"{m['bottom_genre_roi']:.2f}x median ROI", ROSE),
        ]),
        ("Chart 3 of 5", "Release timing is a revenue lever", "Average revenue varies by release month; the peak month in the current data is June.", charts["month_revenue"], [
            ("Best month", MONTHS[m["best_month"] - 1], money(m["best_month_rev"]), TEAL),
            ("Worst month", MONTHS[m["worst_month"] - 1], money(m["worst_month_rev"]), ROSE),
        ]),
        ("Chart 4 of 5", "Revenue and ROI are different goals", "Large budgets lift expected revenue, but lower-budget films can preserve stronger multipliers.", charts["budget_roi"], [
            ("<$10M median", f"{m['low_roi']:.2f}x", "low-budget films", GREEN),
            (">$150M median", f"{m['blockbuster_roi']:.2f}x", "blockbuster films", AMBER),
        ]),
        ("Chart 5 of 5", "Ratings differ clearly by genre", "Ratings are aggregated from the rate table and converted to the 0-10 scale used by the dashboard.", charts["rating_distribution"], [
            ("Highest median", m["rating_top_genre"], f"{m['rating_top']:.2f} / 10", AMBER),
            ("Lowest median", m["rating_bottom_genre"], f"{m['rating_bottom']:.2f} / 10", ROSE),
        ]),
    ]
    for idx, (kicker, title, subtitle, chart_path, cards) in enumerate(chart_slides, start=8):
        s = slide()
        add_title(s, kicker, title, subtitle, idx)
        add_image_fit(s, chart_path, 0.75, 1.55, 8.6, 4.35)
        add_card(s, 9.65, 2.0, 2.6, 1.25, cards[0][0], cards[0][1], cards[0][2], cards[0][3])
        add_card(s, 9.65, 3.65, 2.6, 1.25, cards[1][0], cards[1][1], cards[1][2], cards[1][3])
        add_text(s, "Metric basis: all films in `processed_movies.csv` with complete fields for this chart.", 0.78, 6.55, 8.8, 0.25, 10, MUTED)

    s = slide()
    add_title(s, "Findings", "What We Learned", "The four SDS research questions can now be answered with one consistent metric basis.", 13)
    add_card(s, 0.75, 1.55, 2.65, 1.35, "Q1 · Box office", f"r = {m['budget_revenue_corr']:.2f}", "budget strongly predicts revenue", BLUE)
    add_card(s, 3.65, 1.55, 2.65, 1.35, "Q2 · Genres", f"{m['top_genre']}", f"{m['top_genre_roi']:.2f}x median ROI", PURPLE)
    add_card(s, 6.55, 1.55, 2.65, 1.35, "Q3 · Timing", MONTHS[m["best_month"] - 1], f"{money(m['best_month_rev'])} avg revenue", TEAL)
    add_card(s, 9.45, 1.55, 2.65, 1.35, "Q4 · Ratings", m["rating_sweet_spot"], "highest avg revenue band", AMBER)
    add_bullets(s, [
        f"{pct(m['profit_rate'])} of valid-ROI films are profitable in the current dataset.",
        f"Profitable films score {m['profit_rating_diff']:.2f} rating points higher on average.",
        f"Films longer than 150 minutes post {m['long_roi']:.2f}x median ROI.",
        f"{pct(m['multi_genre_share'])} of all films list two or more genres, so multi-genre filtering matters.",
    ], 1.0, 3.75, 10.8, 1.85, 16)

    s = slide()
    add_title(s, "Model", "Prediction Module Architecture", "Four features in, three independently trained GradientBoosting models out.", 14)
    if "image10" in media:
        add_image_fit(s, media["image10"], 0.8, 1.5, 11.7, 3.0)
    add_card(s, 1.05, 5.25, 2.2, 0.9, "Estimators", "300", "trees per model", BLUE)
    add_card(s, 3.6, 5.25, 2.2, 0.9, "Learning rate", "0.05", "shrinkage factor", TEAL)
    add_card(s, 6.15, 5.25, 2.2, 0.9, "Max depth", "4", "guards overfit", PURPLE)
    add_card(s, 8.7, 5.25, 2.2, 0.9, "Subsample", "0.8", "stochastic boosting", AMBER)

    pred = predict_demo(m["df"])
    s = slide()
    add_title(s, "Demo", "Predictor in Action", "A hypothetical action film returns forecasts plus real comparable films from the dataset.", 15)
    add_card(s, 0.8, 1.45, 2.3, 1.05, "Input", "$80M", "Action · July · 125 min", BLUE)
    add_card(s, 3.35, 1.45, 2.3, 1.05, "Predicted revenue", money(pred["predicted_revenue"]), f"{money(pred['lower_revenue'])} - {money(pred['upper_revenue'])}", TEAL)
    add_card(s, 5.9, 1.45, 2.3, 1.05, "Predicted ROI", f"{pred['predicted_roi']:.2f}x", "GradientBoostingRegressor", PURPLE)
    add_card(s, 8.45, 1.45, 2.3, 1.05, "Profit probability", pct(pred["profit_probability"]), "P(revenue > budget)", GREEN)
    add_table(s, pred["comparables"], 0.8, 3.05, 11.7, 2.85)
    add_text(s, "Comparable films are real rows selected by nearest budget within the same primary genre.", 0.85, 6.25, 10.7, 0.28, 11, MUTED)

    s = slide()
    add_text(s, "26DPWPROJECT · END OF PRESENTATION", 0.75, 0.65, 5.5, 0.25, 10, BLUE, True)
    add_text(s, "Thank you.", 0.75, 1.45, 5.5, 0.8, 44, TEXT, True)
    add_text(s, "Questions, feedback, and live demo welcomed.", 0.78, 2.35, 6.2, 0.35, 18, TEAL, True)
    add_card(s, 0.85, 3.55, 4.0, 1.2, "Repository", "github.com/CandaceZcc/26dpwProject", "run_windows.bat or ./run_unix.sh", BLUE)
    add_card(s, 5.2, 3.55, 4.0, 1.2, "Project scope", "Database · Dashboard · Model", "SQL parsing · charts · statistics", PURPLE)
    add_text(s, "Team: ZHOU Can · CHEN Ziming · DAI Ling · LIAO Shuaiyu · PAN Meihao · ZHU Junyu", 0.85, 5.55, 11.0, 0.35, 14, MUTED)

    try:
        prs.save(OUT)
    except PermissionError:
        prs.save(FALLBACK_OUT)


def add_table(slide, df: pd.DataFrame, x: float, y: float, w: float, h: float) -> None:
    rows, cols = len(df) + 1, len(df.columns)
    table_shape = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    widths = [3.1, 0.8, 1.25, 1.25, 0.85, 0.85, 1.25]
    for i, width in enumerate(widths[:cols]):
        table.columns[i].width = Inches(width)
    for c, col in enumerate(df.columns):
        cell = table.cell(0, c)
        cell.text = col
        style_cell(cell, fill=PANEL, color=MUTED, bold=True, size=8)
    for r in range(len(df)):
        for c, col in enumerate(df.columns):
            cell = table.cell(r + 1, c)
            cell.text = str(df.iloc[r, c])
            style_cell(cell, fill=CARD if r % 2 == 0 else BG, color=TEXT, bold=False, size=8)


def style_cell(cell, fill: str, color: str, bold: bool, size: int) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb(fill)
    for p in cell.text_frame.paragraphs:
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = rgb(color)
        p.alignment = PP_ALIGN.LEFT


def predict_demo(df: pd.DataFrame) -> dict:
    from app.predict import _comparable_table, predict_movie

    budget = 80_000_000
    genre = "Action"
    month = 7
    runtime = 125
    result = predict_movie(df, budget, genre, month, runtime)
    comps = _comparable_table(df, budget, genre)
    return {**result, "comparables": comps}


def main() -> None:
    random.seed(42)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    df = load_df()
    m = compute_metrics(df)
    write_metric_check(m)
    media = extract_old_media()
    charts = {
        "budget_revenue": chart_budget_revenue(m),
        "genre_roi": chart_genre_roi(m),
        "month_revenue": chart_month_revenue(m),
        "budget_roi": chart_budget_roi(m),
        "rating_distribution": chart_rating_distribution(m),
    }
    build_deck(m, charts, media)
    print(f"wrote {OUT}")
    print(f"wrote {METRIC_MD}")


if __name__ == "__main__":
    main()
