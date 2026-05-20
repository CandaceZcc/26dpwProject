from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.update_ppt_speaker4 import (
    BACKUP,
    prepare_crops,
    update_talk_map,
    update_speaker3_text_and_images,
    rebuild_slide_11,
    rebuild_slide_12,
    rebuild_slide_13,
    rebuild_slide_14,
    rebuild_slide_15,
    rebuild_slide_16,
    replace_ui_screenshots,
    global_text_cleanup,
)


OUT_DIR = BACKUP.parent / "data" / "ppt_debug_variants"
OUT_DIR.mkdir(parents=True, exist_ok=True)
crops = prepare_crops()


def save_variant(name: str, ops):
    prs = Presentation(BACKUP)
    for op in ops:
        op(prs)
    global_text_cleanup(prs)
    out = OUT_DIR / f"{name}.pptx"
    prs.save(out)
    print(out)


save_variant("only_global_text", [])
save_variant("talk_map", [lambda prs: update_talk_map(prs.slides[1])])
save_variant("speaker3_images", [lambda prs: update_speaker3_text_and_images(prs.slides[7], prs.slides[8], prs.slides[9], crops)])
save_variant("speaker4_11", [lambda prs: rebuild_slide_11(prs.slides[10], crops)])
save_variant("speaker4_12", [lambda prs: rebuild_slide_12(prs.slides[11])])
save_variant("speaker4_13", [lambda prs: rebuild_slide_13(prs.slides[12], crops)])
save_variant("speaker4_14", [lambda prs: rebuild_slide_14(prs.slides[13], crops)])
save_variant("speaker4_15", [lambda prs: rebuild_slide_15(prs.slides[14], crops)])
save_variant("speaker4_16", [lambda prs: rebuild_slide_16(prs.slides[15])])
save_variant("ui_screenshots", [lambda prs: replace_ui_screenshots(prs, crops)])
