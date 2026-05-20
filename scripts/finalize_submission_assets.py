from __future__ import annotations

import re
import shutil
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


ROOT = Path(__file__).resolve().parents[1]
PPTX = ROOT / "DPW_PPT_Final.pptx"
PPTX_BACKUP = ROOT / "DPW_PPT_Final.before_final_cleanup.pptx"


def iter_text_shapes(slide):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            yield shape


def set_shape_text(shape, text: str) -> None:
    shape.text_frame.clear()
    p = shape.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text


def replace_in_text_frame(shape, replacements: list[tuple[str, str]]) -> None:
    text = shape.text
    new_text = text
    for pattern, replacement in replacements:
        new_text = re.sub(pattern, replacement, new_text)
    if new_text != text:
        set_shape_text(shape, new_text)


def remove_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def add_or_replace_page_marker(slide, slide_no: int, total: int) -> None:
    marker = f"{slide_no:02d} / {total:02d}"
    for shape in iter_text_shapes(slide):
        if shape.name == "final-page-marker":
            set_shape_text(shape, marker)
            return

    box = slide.shapes.add_textbox(Inches(11.35), Inches(6.94), Inches(0.72), Inches(0.22))
    box.name = "final-page-marker"
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = marker
    r.font.size = Pt(7.5)
    r.font.name = "Aptos"
    r.font.color.rgb = RGBColor(120, 128, 140)


def finalize_pptx() -> None:
    if not PPTX_BACKUP.exists():
        shutil.copy2(PPTX, PPTX_BACKUP)

    prs = Presentation(PPTX)
    total = len(prs.slides)

    replacements = [
        (r"\b(\d{1,2})\s*/\s*(18|23)\b", lambda m: f"{int(m.group(1)):02d} / {total:02d}"),
        (r"\bfliter\b", "filter"),
        (r"\bdum\s*ps\b", "dumps"),
        (r"\b3 machine learning\b", "3 model outputs"),
        (r"\b2 statistical test\b", "2 statistical test groups"),
        (r"Gradient Boosting", "Gradient Boosting"),
        (r"scale stabilisation", "scale stabilization"),
        (r"polarising", "polarizing"),
        (r"statistical significance tests, machine learning predictions and exportable insights", "statistical significance tests, machine-learning predictions, and exportable insights"),
    ]

    # Slide 2 contained an old Talk Map layer plus the final version. Remove the old
    # layer by matching text that only appears in the stale block.
    slide2 = prs.slides[1]
    stale_markers = (
        "Presentation Flow by Team Responsibility",
        "The sequence follows the dashboard sidebar",
        "Problem framing, source data, cleaning",
        "Demo & Q&A: 5 min",
    )
    for shape in list(iter_text_shapes(slide2)):
        if any(marker in shape.text for marker in stale_markers):
            remove_shape(shape)

    for slide_no, slide in enumerate(prs.slides, start=1):
        for shape in iter_text_shapes(slide):
            replace_in_text_frame(shape, replacements)
            if re.search(r"\b\d{2}\s*/\s*22\b", shape.text):
                set_shape_text(shape, re.sub(r"\b\d{2}\s*/\s*22\b", f"{slide_no:02d} / {total:02d}", shape.text))
        add_or_replace_page_marker(slide, slide_no, total)

    prs.save(PPTX)


def audit_pptx_text() -> None:
    prs = Presentation(PPTX)
    total = len(prs.slides)
    problems: list[str] = []
    all_text = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_text = "\n".join(shape.text for shape in iter_text_shapes(slide))
        all_text.append(slide_text)
        if re.search(r"/\s*(18|23)\b", slide_text):
            problems.append(f"slide {i}: stale page total found")
        if f"{i:02d} / {total:02d}" not in slide_text:
            problems.append(f"slide {i}: final page marker missing")
    joined = "\n".join(all_text)
    for token in ["fliter", "Presentation Flow by Team Responsibility"]:
        if token in joined:
            problems.append(f"stale token still present: {token}")
    if re.search(r"\bdum\s+ps\b", joined):
            problems.append(f"stale token still present: {token}")
    if problems:
        raise SystemExit("\n".join(problems))


if __name__ == "__main__":
    finalize_pptx()
    audit_pptx_text()
    with zipfile.ZipFile(PPTX) as zf:
        slides = [n for n in zf.namelist() if re.match(r"ppt/slides/slide\d+\.xml$", n)]
    print(f"Finalized {PPTX.name}: {len(slides)} slides, page markers 01-{len(slides):02d}.")
