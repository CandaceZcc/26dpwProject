from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

from scripts.revise_dpw_ppt import (
    AMBER,
    ASSET_DIR,
    BG,
    BLUE,
    BORDER,
    CARD,
    GREEN,
    MONTHS,
    MUTED,
    OUT,
    PANEL,
    PURPLE,
    ROSE,
    TEAL,
    TEXT,
    add_bullets,
    add_card,
    add_image_fit,
    add_table,
    add_text,
    chart_budget_revenue,
    chart_budget_roi,
    chart_genre_roi,
    chart_month_revenue,
    chart_rating_distribution,
    compute_metrics,
    extract_old_media,
    load_df,
    money,
    pct,
    predict_demo,
    rgb,
    set_bg,
    write_metric_check,
)


TOTAL = 18
SPEAKER_OUT = OUT.parent / "DPW_PPT_speaker_structured.pptx"


def add_title(slide, section: str, title: str, subtitle: str, num: int) -> None:
    add_text(slide, section.upper(), 0.55, 0.32, 4.2, 0.24, 9, BLUE, True)
    add_text(slide, f"{num:02d} / {TOTAL}", 11.6, 0.32, 1.1, 0.24, 9, MUTED, True, PP_ALIGN.RIGHT)
    add_text(slide, title, 0.55, 0.64, 8.8, 0.58, 27, TEXT, True)
    if subtitle:
        add_text(slide, subtitle, 0.58, 1.14, 10.6, 0.36, 12, MUTED)


def speaker_footer(slide, speaker: str, time: str, role: str) -> None:
    shape = slide.shapes.add_shape(1, Inches(0.55), Inches(6.88), Inches(12.25), Inches(0.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(PANEL)
    shape.line.color.rgb = rgb(BORDER)
    shape.line.width = Pt(0.6)
    add_text(slide, f"{speaker} · {time}", 0.75, 6.96, 2.6, 0.16, 8, TEAL, True)
    add_text(slide, role, 3.0, 6.96, 9.4, 0.16, 8, MUTED)


def pill(slide, text: str, x: float, y: float, w: float, color: str = BLUE) -> None:
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(CARD)
    shape.line.color.rgb = rgb(color)
    shape.line.width = Pt(1)
    add_text(slide, text, x + 0.1, y + 0.085, w - 0.2, 0.14, 8, color, True, PP_ALIGN.CENTER)


def build_deck(metrics: dict, charts: dict[str, Path], media: dict[str, Path]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def slide():
        s = prs.slides.add_slide(blank)
        set_bg(s)
        return s

    s = slide()
    add_text(s, "26DPWPROJECT · FINAL PRESENTATION", 0.75, 0.55, 5.8, 0.25, 10, BLUE, True)
    add_text(s, "TMDB Movie Analytics\nDashboard", 0.75, 1.25, 7.5, 1.55, 39, TEXT, True)
    add_text(s, "A 15-minute team story: data source → database → dashboard → statistics → prediction → demo.", 0.78, 3.14, 9.7, 0.36, 16, TEAL, True)
    add_card(s, 0.8, 4.35, 2.35, 1.15, "Dataset", f"{metrics['movie_count']:,}", "analysis-ready films", TEAL)
    add_card(s, 3.45, 4.35, 2.35, 1.15, "Financial rows", f"{metrics['valid_roi_count']:,}", "budget + revenue", BLUE)
    add_card(s, 6.1, 4.35, 2.35, 1.15, "Views", "8", "dashboard sections", PURPLE)
    add_card(s, 8.75, 4.35, 2.35, 1.15, "Demo", "5 min", "live system walkthrough", AMBER)
    add_text(s, "ZHOU Can · CHEN Ziming · DAI Ling · LIAO Shuaiyu · PAN Meihao · ZHU Junyu", 0.8, 6.22, 9.8, 0.26, 11, MUTED)

    s = slide()
    add_title(s, "Talk map", "按演讲分工重构后的汇报路线", "顺序跟随 dashboard 侧边栏视图，同时让每位成员有清晰负责段落。", 2)
    rows = [
        ("1", "Project + Data Cleaning", "2.5 min", "背景、数据来源、Python 清洗"),
        ("2", "Database + ER", "2.5 min", "关系设计、核心表、SQL 导出"),
        ("3", "Dashboard Views", "3 min", "Overview / Revenue / Genres / Time"),
        ("4", "Chi-Square", "2 min", "分类变量显著性检验"),
        ("5", "Regression + Prediction", "2.5 min", "线性回归与 Gradient Boosting"),
        ("6", "Insights + Export + Summary", "2.5 min", "自动洞察、导出、局限与未来方向"),
    ]
    x_positions = [0.75, 2.75, 4.95, 7.15, 9.35, 11.0]
    colors = [BLUE, TEAL, PURPLE, AMBER, GREEN, ROSE]
    for i, row in enumerate(rows):
        x = x_positions[i]
        add_card(s, x, 1.75, 1.75, 3.55, row[0], row[1], f"{row[2]}\n{row[3]}", colors[i])
    add_text(s, "Demo & Q&A: 5 min · 第三位或第六位主操作，其他成员辅助回答。", 0.9, 5.95, 10.8, 0.35, 16, TEXT, True, PP_ALIGN.CENTER)

    s = slide()
    add_title(s, "Speaker 1", "项目背景：从分散电影数据到可交互决策系统", "目标是帮助分析票房规律、类型盈利、档期选择，并实现基础票房预测。", 3)
    speaker_footer(s, "第一位", "约 2.5 min", "项目介绍 / 问题背景 / 数据清洗")
    add_bullets(s, [
        "电影行业数据量大但分散，投资决策和档期选择需要数据支撑。",
        "系统目标：把 TMDB SQL 数据转成可筛选、可视化、可预测的 dashboard。",
        "核心问题：预算是否推高票房？哪些类型更赚钱？什么时候上映更好？评分如何影响表现？",
    ], 0.8, 1.85, 6.8, 2.7, 17)
    add_card(s, 8.1, 1.9, 2.6, 1.15, "Raw source", "TMDB SQL", "movie, genres, rate, credit...", BLUE)
    add_card(s, 8.1, 3.35, 2.6, 1.15, "Final table", f"{metrics['movie_count']:,}", "processed movies", TEAL)
    add_card(s, 8.1, 4.8, 2.6, 1.15, "Complete finance", f"{metrics['valid_roi_count']:,}", "budget / revenue rows", PURPLE)

    s = slide()
    add_title(s, "Speaker 1", "Python 清洗流程：无需本地数据库", "`scripts/build_dataset.py` 直接解析 SQL dump，输出 dashboard 可读 CSV。", 4)
    speaker_footer(s, "第一位", "约 2.5 min", "数据清洗流程与字段构造")
    add_bullets(s, [
        "Direct parse: 读取 `.sql` 文件中的 INSERT statements，不要求用户安装 MySQL。",
        "Join logic: movie + genres + link_genres + rate，形成一张分析宽表。",
        "Feature build: year / release_month / ROI / avg_rating / primary_genre / vote_count。",
        "Rating scale: 0.5-5 分制转换为 0-10 分制，匹配 dashboard 展示。",
        "Outlier handling: ROI 图表和训练中使用分位数截断，避免极端值压缩主体分布。",
    ], 0.85, 1.7, 8.6, 3.8, 15)
    add_card(s, 9.8, 1.9, 2.3, 1.0, "Output", "CSV", "`data/processed_movies.csv`", GREEN)
    add_card(s, 9.8, 3.15, 2.3, 1.0, "Fields", "17", "analysis columns", BLUE)
    add_card(s, 9.8, 4.4, 2.3, 1.0, "Skipped", "87", "missing release date", ROSE)

    s = slide()
    add_title(s, "Speaker 2", "数据库设计：关系型 schema 支撑可复现分析", "核心设计遵循规范化思路，再通过 Python 聚合成 dashboard 数据集。", 5)
    speaker_footer(s, "第二位", "约 2.5 min", "数据库设计 / ER 图 / SQL 导出")
    if "image1" in media:
        add_image_fit(s, media["image1"], 0.75, 1.45, 7.6, 3.95)
    add_card(s, 8.75, 1.55, 2.9, 1.0, "Core relation", "movie ↔ genres", "many-to-many via link_genres", BLUE)
    add_card(s, 8.75, 2.85, 2.9, 1.0, "Rating relation", "movie ↔ rate", "one movie, many user ratings", TEAL)
    add_card(s, 8.75, 4.15, 2.9, 1.0, "People relation", "cast / credit", "actors, directors, writers", PURPLE)
    add_text(s, "项目目录当前包含 23 个 SQL dump files；PPT 中统一按实际提交包口径表述。", 0.82, 5.95, 10.5, 0.28, 10, MUTED)

    s = slide()
    add_title(s, "Speaker 2", "核心表如何映射到分析问题", "ER 图讲结构，这页讲每张表如何服务 dashboard。", 6)
    speaker_footer(s, "第二位", "约 2.5 min", "核心表说明与 SQL 导出")
    table_df = __import__("pandas").DataFrame([
        ["movie", "title, budget, revenue, runtime, release_date", "票房、预算、时长、月份分析"],
        ["genres", "genre dictionary", "类型名称标准化"],
        ["link_genres", "movie-genre relation", "多类型展开与 primary_genre"],
        ["rate", "user_id, tmdbid, rating, timestamp", "评分分布与口碑分析"],
        ["cast / credit", "people and roles", "后续可扩展导演/演员特征"],
        ["production_companies", "company metadata", "后续可扩展公司影响分析"],
    ], columns=["Table", "Key fields", "Used for"])
    add_table(s, table_df, 0.75, 1.65, 11.8, 3.9)
    add_text(s, "SQL files are exported separately, then parsed by Python so the app can run without a local database server.", 0.9, 6.0, 10.8, 0.35, 14, TEAL, True, PP_ALIGN.CENTER)

    s = slide()
    add_title(s, "Speaker 3", "系统整体界面：一个入口，八个分析视图", "这一段按侧边栏顺序讲 Overview / Revenue / Genres / Time。", 7)
    speaker_footer(s, "第三位", "约 3 min", "Overview / Revenue / Genres / Time")
    if "image2" in media:
        add_image_fit(s, media["image2"], 0.75, 1.45, 8.25, 4.65)
    add_card(s, 9.35, 1.6, 2.55, 1.0, "Filter rail", "Shared state", "year, genre, min votes", BLUE)
    add_card(s, 9.35, 2.85, 2.55, 1.0, "KPI row", "6 cards", "movies, revenue, ROI, rating", TEAL)
    add_card(s, 9.35, 4.1, 2.55, 1.0, "Views", "8", "Overview to Export", PURPLE)
    add_text(s, "讲法：先说明所有筛选会同步刷新 KPI、图表和导出结果。", 0.85, 6.28, 10.4, 0.28, 11, MUTED)

    s = slide()
    add_title(s, "Speaker 3", "Revenue 视图：预算推高票房，但不保证 ROI", "先用预算-票房回答商业规模，再用预算-ROI 解释回报率风险。", 8)
    speaker_footer(s, "第三位", "约 3 min", "Revenue view")
    add_image_fit(s, charts["budget_revenue"], 0.75, 1.55, 5.85, 3.65)
    add_image_fit(s, charts["budget_roi"], 6.85, 1.55, 5.85, 3.65)
    add_card(s, 1.05, 5.65, 2.5, 0.9, "Budget → Revenue", f"r = {metrics['budget_revenue_corr']:.2f}", "clear positive correlation", BLUE)
    add_card(s, 3.9, 5.65, 2.5, 0.9, "Low budget ROI", f"{metrics['low_roi']:.2f}x", "< $10M median", GREEN)
    add_card(s, 6.75, 5.65, 2.5, 0.9, "Blockbuster ROI", f"{metrics['blockbuster_roi']:.2f}x", "> $150M median", AMBER)
    add_card(s, 9.6, 5.65, 2.5, 0.9, "Conclusion", "Scale ≠ efficiency", "two goals differ", PURPLE)

    s = slide()
    add_title(s, "Speaker 3", "Genres 与 Time：类型和档期决定盈利侧重点", "把 dashboard 的类型排序、评分分布和上映月份曲线合并成一段讲。", 9)
    speaker_footer(s, "第三位", "约 3 min", "Genres view + Time view")
    add_image_fit(s, charts["genre_roi"], 0.65, 1.52, 4.05, 3.25)
    add_image_fit(s, charts["rating_distribution"], 4.82, 1.52, 4.05, 3.25)
    add_image_fit(s, charts["month_revenue"], 8.98, 1.52, 3.7, 3.25)
    add_card(s, 0.95, 5.25, 2.3, 0.9, "Top ROI genre", metrics["top_genre"], f"{metrics['top_genre_roi']:.2f}x median", PURPLE)
    add_card(s, 3.55, 5.25, 2.3, 0.9, "Best month", MONTHS[metrics["best_month"] - 1], money(metrics["best_month_rev"]), TEAL)
    add_card(s, 6.15, 5.25, 2.3, 0.9, "Best rating band", metrics["rating_sweet_spot"], "highest avg revenue", AMBER)
    add_card(s, 8.75, 5.25, 2.3, 0.9, "Interaction", "Metric toggle", "ROI or avg revenue", BLUE)

    s = slide()
    add_title(s, "Speaker 4", "Chi-Square：把直觉结论变成显著性检验", "检验两个分类变量之间是否存在统计上的显著关联。", 10)
    speaker_footer(s, "第四位", "约 2 min", "Chi-Square 卡方检验")
    if "image8" in media:
        add_image_fit(s, media["image8"], 0.75, 1.45, 6.6, 3.65)
    add_bullets(s, [
        "H0: 两个分类变量独立；当 p < 0.05 时拒绝 H0。",
        "系统展示 χ² statistic、p-value、degrees of freedom 和 Cramér's V。",
        "Residual heatmap 显示哪些组合高于/低于独立性假设下的期望频数。",
    ], 7.75, 1.7, 4.5, 2.6, 15)
    add_card(s, 1.0, 5.55, 2.55, 0.85, "Pairing 1", "Genre × Profit", "类型和盈利性", BLUE)
    add_card(s, 3.85, 5.55, 2.55, 0.85, "Pairing 2", "Budget × Profit", "预算等级和盈利", TEAL)
    add_card(s, 6.7, 5.55, 2.55, 0.85, "Pairing 3", "Season × Revenue", "档期和高票房", PURPLE)
    add_card(s, 9.55, 5.55, 2.55, 0.85, "Pairing 4", "Rating × Profit", "评分等级和盈利", AMBER)

    s = slide()
    add_title(s, "Speaker 4", "派生字段让统计检验可解释", "把连续变量转成类别变量，才能进入卡方检验和交叉表热力图。", 11)
    speaker_footer(s, "第四位", "约 2 min", "派生字段与解释方式")
    add_card(s, 0.85, 1.7, 3.2, 1.25, "budget_tier", "<$10M / $10-50M / $50-150M / >$150M", "低/中/高/大片预算", BLUE)
    add_card(s, 4.35, 1.7, 3.2, 1.25, "season", "Winter / Spring / Summer / Fall / Awards", "月份映射为档期", TEAL)
    add_card(s, 7.85, 1.7, 3.2, 1.25, "profitable", "revenue > budget", "盈利 / 未盈利", GREEN)
    add_bullets(s, [
        "讲图时先读 p-value：是否显著。",
        "再读 Cramér's V：关联强度大小。",
        "最后读 residual heatmap：哪些类别组合贡献了差异。",
    ], 1.05, 3.75, 10.2, 1.65, 18)
    add_text(s, "这一页的作用：帮助听众理解系统里的统计模块不是单纯画热力图，而是在检验变量关系。", 1.0, 5.85, 10.9, 0.35, 14, MUTED, False, PP_ALIGN.CENTER)

    s = slide()
    add_title(s, "Speaker 5", "Regression：交互式探索特征和票房关系", "线性回归用于解释关系，预测模块用于输出可操作估计。", 12)
    speaker_footer(s, "第五位", "约 2.5 min", "Regression 线性回归")
    if "image9" in media:
        add_image_fit(s, media["image9"], 0.75, 1.45, 7.2, 3.8)
    add_card(s, 8.35, 1.7, 2.75, 0.9, "Variables", "8", "budget, revenue, ROI...", BLUE)
    add_card(s, 8.35, 2.9, 2.75, 0.9, "Split", "80/20", "train / test", TEAL)
    add_card(s, 8.35, 4.1, 2.75, 0.9, "Output", "R² + line", "fit quality visible", PURPLE)
    add_text(s, "讲法：Regression view 是解释工具，让老师看到任意 X/Y 组合都能快速测试线性关系。", 0.85, 5.85, 10.8, 0.32, 13, MUTED)

    s = slide()
    add_title(s, "Speaker 5", "Prediction：Gradient Boosting 输出票房、ROI 和盈利概率", "四个输入特征，三个模型，最后给出 Top 5 可比较历史影片。", 13)
    speaker_footer(s, "第五位", "约 2.5 min", "机器学习预测模块")
    if "image10" in media:
        add_image_fit(s, media["image10"], 0.85, 1.55, 6.7, 2.9)
    add_card(s, 8.05, 1.55, 2.4, 0.9, "Feature 1", "log(budget)", "scale stabilisation", BLUE)
    add_card(s, 10.65, 1.55, 2.0, 0.9, "Feature 2", "genre", "OrdinalEncoder", TEAL)
    add_card(s, 8.05, 2.75, 2.4, 0.9, "Feature 3", "month", "release timing", PURPLE)
    add_card(s, 10.65, 2.75, 2.0, 0.9, "Feature 4", "runtime", "minutes", AMBER)
    add_card(s, 1.0, 5.05, 2.4, 0.9, "Revenue", "GB Regressor", "log revenue target", BLUE)
    add_card(s, 3.75, 5.05, 2.4, 0.9, "ROI", "GB Regressor", "97th percentile cap", PURPLE)
    add_card(s, 6.5, 5.05, 2.4, 0.9, "Profit", "GB Classifier", "revenue > budget", GREEN)
    add_card(s, 9.25, 5.05, 2.4, 0.9, "Params", "300 · 0.05 · 4", "trees · lr · depth", AMBER)

    pred = predict_demo(metrics["df"])
    s = slide()
    add_title(s, "Speaker 5", "Prediction 输出：把模型结果讲成产品功能", "用一个 Action / July / $80M / 125min 案例展示预测和真实可比影片。", 14)
    speaker_footer(s, "第五位", "约 2.5 min", "预测输出与可比影片")
    add_card(s, 0.8, 1.45, 2.25, 0.98, "Input", "$80M", "Action · July · 125 min", BLUE)
    add_card(s, 3.35, 1.45, 2.25, 0.98, "Revenue", money(pred["predicted_revenue"]), f"{money(pred['lower_revenue'])} - {money(pred['upper_revenue'])}", TEAL)
    add_card(s, 5.9, 1.45, 2.25, 0.98, "ROI", f"{pred['predicted_roi']:.2f}x", "predicted multiplier", PURPLE)
    add_card(s, 8.45, 1.45, 2.25, 0.98, "Profit prob.", pct(pred["profit_probability"]), "P(revenue > budget)", GREEN)
    add_table(s, pred["comparables"], 0.8, 3.0, 11.7, 2.75)
    add_text(s, "Comparable rows are real films selected by nearest budget within the same primary genre.", 0.85, 6.12, 10.8, 0.28, 10, MUTED)

    s = slide()
    add_title(s, "Speaker 6", "Insights：把图表结果自动总结成可讲结论", "最后一位负责把 dashboard 的洞察、导出能力和系统边界收束起来。", 15)
    speaker_footer(s, "第六位", "约 2.5 min", "Insights / Export / 总结")
    add_card(s, 0.75, 1.55, 2.65, 1.15, "Highest ROI", metrics["top_genre"], f"{metrics['top_genre_roi']:.2f}x median", PURPLE)
    add_card(s, 3.65, 1.55, 2.65, 1.15, "Best month", MONTHS[metrics["best_month"] - 1], money(metrics["best_month_rev"]), TEAL)
    add_card(s, 6.55, 1.55, 2.65, 1.15, "Profit rate", pct(metrics["profit_rate"]), "valid ROI films", GREEN)
    add_card(s, 9.45, 1.55, 2.65, 1.15, "Long films", f"{metrics['long_roi']:.2f}x", ">150 min median ROI", AMBER)
    add_bullets(s, [
        f"Multi-genre strategy: {pct(metrics['multi_genre_share'])} of movies include two or more genres.",
        f"Rating sweet spot: {metrics['rating_sweet_spot']} rating band has the highest average revenue.",
        "Insights view turns raw charts into presenter-friendly cards, trends and findings.",
    ], 1.0, 3.65, 10.8, 1.6, 17)

    s = slide()
    add_title(s, "Speaker 6", "Export 与总结：系统亮点、局限和未来方向", "这页作为正式讲解收尾，后面进入 5 分钟 Demo。", 16)
    speaker_footer(s, "第六位", "约 2.5 min", "Export / limitations / future work")
    add_card(s, 0.8, 1.55, 3.2, 1.2, "Export", "CSV download", "current filtered subset", BLUE)
    add_card(s, 4.35, 1.55, 3.2, 1.2, "Strengths", "No DB required", "one-click launch + 8 views", GREEN)
    add_card(s, 7.9, 1.55, 3.2, 1.2, "Limitations", "Wide interval", "4-feature prediction only", ROSE)
    add_bullets(s, [
        "System highlight: database design, reproducible parsing, interactive charts, statistics and ML prediction form a full pipeline.",
        "Limitation: prediction uncertainty is wide and current model does not use director, actor or production-company effects.",
        "Future work: connect real-time TMDB API, add people/company features, compare more statistical and ML models.",
    ], 1.0, 3.45, 10.7, 2.0, 16)

    s = slide()
    add_title(s, "Demo", "5-minute live demo route", "建议第三位或第六位操作，其他成员辅助回答具体模块问题。", 17)
    pill(s, "0-1 min · Open + filters", 0.9, 1.65, 2.6, BLUE)
    pill(s, "1-2 min · Revenue / Genres", 3.75, 1.65, 2.8, TEAL)
    pill(s, "2-3.5 min · Prediction", 6.8, 1.65, 2.6, PURPLE)
    pill(s, "3.5-5 min · Q&A", 9.65, 1.65, 2.3, AMBER)
    add_bullets(s, [
        "Open Streamlit and set filters, e.g. 2010-2023, Action, min votes 1000.",
        "Switch Revenue and Genres views; point out budget correlation, ROI and genre ranking.",
        "Enter prediction case; show revenue range, ROI, profit probability and comparable films.",
        "Use Q&A to let each member answer questions about their own module.",
    ], 1.0, 2.75, 10.8, 2.6, 18)

    s = slide()
    add_text(s, "26DPWPROJECT · END OF PRESENTATION", 0.75, 0.65, 5.5, 0.25, 10, BLUE, True)
    add_text(s, "Thank you.", 0.75, 1.45, 5.5, 0.8, 44, TEXT, True)
    add_text(s, "Questions, feedback, and live demo welcomed.", 0.78, 2.35, 6.2, 0.35, 18, TEAL, True)
    add_card(s, 0.85, 3.55, 4.0, 1.2, "Repository", "github.com/CandaceZcc/26dpwProject", "run_windows.bat or ./run_unix.sh", BLUE)
    add_card(s, 5.2, 3.55, 4.0, 1.2, "Final structure", "6 speakers + demo", "aligned with sidebar views", PURPLE)
    add_text(s, "Team: ZHOU Can · CHEN Ziming · DAI Ling · LIAO Shuaiyu · PAN Meihao · ZHU Junyu", 0.85, 5.55, 11.0, 0.35, 14, MUTED)

    prs.save(SPEAKER_OUT)


def main() -> None:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    metrics = compute_metrics(load_df())
    write_metric_check(metrics)
    media = extract_old_media()
    charts = {
        "budget_revenue": chart_budget_revenue(metrics),
        "budget_roi": chart_budget_roi(metrics),
        "genre_roi": chart_genre_roi(metrics),
        "month_revenue": chart_month_revenue(metrics),
        "rating_distribution": chart_rating_distribution(metrics),
    }
    build_deck(metrics, charts, media)
    print(f"wrote {SPEAKER_OUT}")


if __name__ == "__main__":
    main()
